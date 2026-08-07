from __future__ import annotations

import csv
import hashlib
import io
import json
from dataclasses import asdict, dataclass, replace
from datetime import UTC, datetime
from decimal import Decimal
from enum import StrEnum
from typing import Any
from uuid import uuid4

from openpyxl import Workbook
from pptx import Presentation
from reportlab.lib.pagesizes import A4
from reportlab.pdfgen.canvas import Canvas


class ReportType(StrEnum):
    MANAGEMENT_PACK = "management_pack"
    BOARD_RISK_PACK = "board_risk_pack"
    FORECAST_REPORT = "forecast_report"
    LAGEBERICHT_DRAFT = "lagebericht_draft"
    AUDIT_EVIDENCE_PACK = "audit_evidence_pack"


class ReportStatus(StrEnum):
    DRAFT = "draft"
    APPROVED = "approved"


class ExportFormat(StrEnum):
    JSON = "json"
    CSV = "csv"
    EXCEL = "xlsx"
    PDF = "pdf"
    POWERPOINT = "pptx"


@dataclass(frozen=True, slots=True)
class ReportValue:
    key: str
    value: Decimal | str | int | float | bool | None
    unit: str | None
    snapshot_id: str
    run_id: str
    run_status: str

    @property
    def is_approved(self) -> bool:
        return self.run_status.lower() == "approved"


@dataclass(frozen=True, slots=True)
class NarrativeStatement:
    text: str
    source_refs: tuple[str, ...]


@dataclass(frozen=True, slots=True)
class ReportSection:
    section_id: str
    title: str
    values: tuple[ReportValue, ...] = ()
    statements: tuple[NarrativeStatement, ...] = ()
    metadata: tuple[tuple[str, str], ...] = ()


@dataclass(frozen=True, slots=True)
class ReportTemplate:
    template_id: str
    version: int
    name: str
    report_type: ReportType
    external: bool
    required_sections: tuple[str, ...]


@dataclass(frozen=True, slots=True)
class ReportApproval:
    approver: str
    approved_at: datetime


@dataclass(frozen=True, slots=True)
class ReportArtifact:
    report_id: str
    template_id: str
    template_version: int
    report_type: ReportType
    external: bool
    status: ReportStatus
    generated_at: datetime
    sections: tuple[ReportSection, ...]
    content_hash: str
    approval: ReportApproval | None = None


class TemplateRegistry:
    def __init__(self, templates: tuple[ReportTemplate, ...] = ()) -> None:
        self._templates: dict[tuple[str, int], ReportTemplate] = {}
        for template in templates:
            self.register(template)

    def register(self, template: ReportTemplate) -> None:
        key = (template.template_id, template.version)
        if key in self._templates:
            raise ValueError(f"template already exists: {template.template_id} v{template.version}")
        self._templates[key] = template

    def get(self, template_id: str, version: int) -> ReportTemplate:
        try:
            return self._templates[(template_id, version)]
        except KeyError as exc:
            raise KeyError(f"unknown template: {template_id} v{version}") from exc

    def list(self) -> tuple[ReportTemplate, ...]:
        return tuple(sorted(self._templates.values(), key=lambda item: (item.template_id, item.version)))


class InMemoryReportRepository:
    def __init__(self) -> None:
        self._reports: dict[str, ReportArtifact] = {}

    def save(self, report: ReportArtifact) -> None:
        current = self._reports.get(report.report_id)
        if current is not None and current.status is ReportStatus.APPROVED:
            raise ValueError("approved reports are immutable")
        self._reports[report.report_id] = report

    def get(self, report_id: str) -> ReportArtifact:
        try:
            return self._reports[report_id]
        except KeyError as exc:
            raise KeyError(f"unknown report: {report_id}") from exc

    def list(self) -> tuple[ReportArtifact, ...]:
        return tuple(self._reports.values())


class ReportingFactory:
    def __init__(self, registry: TemplateRegistry, repository: InMemoryReportRepository) -> None:
        self._registry = registry
        self._repository = repository

    def generate(
        self,
        template_id: str,
        template_version: int,
        sections: tuple[ReportSection, ...],
    ) -> ReportArtifact:
        template = self._registry.get(template_id, template_version)
        self._validate_sections(template, sections)
        generated_at = datetime.now(UTC)
        report = ReportArtifact(
            report_id=str(uuid4()),
            template_id=template.template_id,
            template_version=template.version,
            report_type=template.report_type,
            external=template.external,
            status=ReportStatus.DRAFT,
            generated_at=generated_at,
            sections=sections,
            content_hash=self._content_hash(template, sections),
        )
        self._repository.save(report)
        return report

    def approve(self, report_id: str, approver: str) -> ReportArtifact:
        if not approver.strip():
            raise ValueError("approver is required")
        report = self._repository.get(report_id)
        if report.status is ReportStatus.APPROVED:
            return report
        approved = replace(
            report,
            status=ReportStatus.APPROVED,
            approval=ReportApproval(approver=approver.strip(), approved_at=datetime.now(UTC)),
        )
        self._repository.save(approved)
        return approved

    def get(self, report_id: str) -> ReportArtifact:
        return self._repository.get(report_id)

    def list_templates(self) -> tuple[ReportTemplate, ...]:
        return self._registry.list()

    @staticmethod
    def _validate_sections(
        template: ReportTemplate,
        sections: tuple[ReportSection, ...],
    ) -> None:
        section_ids = {section.section_id for section in sections}
        missing = set(template.required_sections) - section_ids
        if missing:
            raise ValueError(f"missing required report sections: {sorted(missing)}")
        for section in sections:
            for value in section.values:
                if not value.snapshot_id or not value.run_id:
                    raise ValueError(f"value {value.key} is missing lineage")
                if not value.is_approved:
                    raise ValueError(f"value {value.key} does not reference an approved run")
            for statement in section.statements:
                if not statement.source_refs:
                    raise ValueError("every material statement requires source references")

    @staticmethod
    def _content_hash(
        template: ReportTemplate,
        sections: tuple[ReportSection, ...],
    ) -> str:
        payload = {
            "template": {
                "id": template.template_id,
                "version": template.version,
                "type": template.report_type.value,
            },
            "sections": _json_ready(sections),
        }
        canonical = json.dumps(payload, sort_keys=True, separators=(",", ":"), ensure_ascii=False)
        return hashlib.sha256(canonical.encode("utf-8")).hexdigest()


class ReportExporter:
    def export(self, report: ReportArtifact, export_format: ExportFormat) -> bytes:
        if report.external and report.status is not ReportStatus.APPROVED:
            raise ValueError("external reports require human approval before export")
        if export_format is ExportFormat.JSON:
            return self._json(report)
        if export_format is ExportFormat.CSV:
            return self._csv(report)
        if export_format is ExportFormat.EXCEL:
            return self._excel(report)
        if export_format is ExportFormat.PDF:
            return self._pdf(report)
        if export_format is ExportFormat.POWERPOINT:
            return self._powerpoint(report)
        raise ValueError(f"unsupported export format: {export_format}")

    @staticmethod
    def _json(report: ReportArtifact) -> bytes:
        return json.dumps(_json_ready(report), ensure_ascii=False, indent=2).encode("utf-8")

    @staticmethod
    def _csv(report: ReportArtifact) -> bytes:
        buffer = io.StringIO()
        writer = csv.writer(buffer)
        writer.writerow(["section", "key", "value", "unit", "snapshot_id", "run_id"])
        for section in report.sections:
            for value in section.values:
                writer.writerow(
                    [
                        section.section_id,
                        value.key,
                        value.value,
                        value.unit or "",
                        value.snapshot_id,
                        value.run_id,
                    ]
                )
        return buffer.getvalue().encode("utf-8")

    @staticmethod
    def _excel(report: ReportArtifact) -> bytes:
        workbook = Workbook()
        default = workbook.active
        workbook.remove(default)
        for section in report.sections:
            sheet = workbook.create_sheet(title=_safe_sheet_name(section.title))
            sheet.append(["Key", "Value", "Unit", "Snapshot", "Run"])
            for value in section.values:
                sheet.append(
                    [
                        value.key,
                        _excel_value(value.value),
                        value.unit or "",
                        value.snapshot_id,
                        value.run_id,
                    ]
                )
            if section.statements:
                sheet.append([])
                sheet.append(["Narrative", "Sources"])
                for statement in section.statements:
                    sheet.append([statement.text, ", ".join(statement.source_refs)])
        buffer = io.BytesIO()
        workbook.save(buffer)
        return buffer.getvalue()

    @staticmethod
    def _pdf(report: ReportArtifact) -> bytes:
        buffer = io.BytesIO()
        canvas = Canvas(buffer, pagesize=A4)
        width, height = A4
        y = height - 50
        canvas.setTitle(f"{report.report_type.value} {report.report_id}")
        canvas.setFont("Helvetica-Bold", 14)
        canvas.drawString(50, y, report.report_type.value.replace("_", " ").title())
        y -= 28
        for section in report.sections:
            if y < 100:
                canvas.showPage()
                y = height - 50
            canvas.setFont("Helvetica-Bold", 11)
            canvas.drawString(50, y, section.title[:90])
            y -= 18
            canvas.setFont("Helvetica", 9)
            for value in section.values:
                if y < 70:
                    canvas.showPage()
                    y = height - 50
                    canvas.setFont("Helvetica", 9)
                text = f"{value.key}: {value.value} {value.unit or ''} | run={value.run_id}"
                canvas.drawString(60, y, text[:115])
                y -= 14
            for statement in section.statements:
                if y < 70:
                    canvas.showPage()
                    y = height - 50
                    canvas.setFont("Helvetica", 9)
                canvas.drawString(60, y, statement.text[:115])
                y -= 14
            y -= 8
        canvas.save()
        return buffer.getvalue()

    @staticmethod
    def _powerpoint(report: ReportArtifact) -> bytes:
        presentation = Presentation()
        title_slide = presentation.slides.add_slide(presentation.slide_layouts[0])
        title_slide.shapes.title.text = report.report_type.value.replace("_", " ").title()
        title_slide.placeholders[1].text = f"Report {report.report_id}"
        for section in report.sections:
            slide = presentation.slides.add_slide(presentation.slide_layouts[1])
            slide.shapes.title.text = section.title
            frame = slide.placeholders[1].text_frame
            frame.clear()
            lines = [f"{value.key}: {value.value} {value.unit or ''}" for value in section.values]
            lines.extend(statement.text for statement in section.statements)
            for index, line in enumerate(lines[:12]):
                paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
                paragraph.text = line
        buffer = io.BytesIO()
        presentation.save(buffer)
        return buffer.getvalue()


def built_in_templates() -> tuple[ReportTemplate, ...]:
    return (
        ReportTemplate(
            template_id="management-pack",
            version=1,
            name="Management Pack",
            report_type=ReportType.MANAGEMENT_PACK,
            external=False,
            required_sections=("kpi", "performance", "forecast", "cash"),
        ),
        ReportTemplate(
            template_id="board-risk-pack",
            version=1,
            name="Board Risk Pack",
            report_type=ReportType.BOARD_RISK_PACK,
            external=False,
            required_sections=("top-risks", "risk-capacity", "stress", "limits"),
        ),
        ReportTemplate(
            template_id="forecast-report",
            version=1,
            name="Forecast Report",
            report_type=ReportType.FORECAST_REPORT,
            external=False,
            required_sections=("assumptions", "distribution", "targets", "model-quality"),
        ),
        ReportTemplate(
            template_id="lagebericht-draft",
            version=1,
            name="Lagebericht Draft",
            report_type=ReportType.LAGEBERICHT_DRAFT,
            external=True,
            required_sections=("economic-report", "forecast", "opportunities", "risks"),
        ),
        ReportTemplate(
            template_id="audit-evidence-pack",
            version=1,
            name="Audit Evidence Pack",
            report_type=ReportType.AUDIT_EVIDENCE_PACK,
            external=False,
            required_sections=("lineage", "approvals", "models", "assumptions", "hashes"),
        ),
    )


def _safe_sheet_name(value: str) -> str:
    cleaned = "".join(char for char in value if char not in "[]:*?/\\")
    return (cleaned or "Report")[:31]


def _excel_value(value: Any) -> Any:
    if isinstance(value, Decimal):
        return float(value)
    return value


def _json_ready(value: Any) -> Any:
    if isinstance(value, Decimal):
        return str(value)
    if isinstance(value, datetime):
        return value.isoformat()
    if isinstance(value, StrEnum):
        return value.value
    if hasattr(value, "__dataclass_fields__"):
        return {key: _json_ready(item) for key, item in asdict(value).items()}
    if isinstance(value, tuple):
        return [_json_ready(item) for item in value]
    if isinstance(value, list):
        return [_json_ready(item) for item in value]
    if isinstance(value, dict):
        return {str(key): _json_ready(item) for key, item in value.items()}
    return value
